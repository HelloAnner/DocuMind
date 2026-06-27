#!/usr/bin/env python3
import json
import mimetypes
import os
import shlex
import subprocess
import tempfile
import time
import urllib.error
import urllib.request
import uuid
import zipfile
from pathlib import Path


BASE_URL = os.environ.get("BASE_URL", "http://123.57.255.204:8089").rstrip("/")
LOGIN_EMAIL = os.environ.get("LOGIN_EMAIL", "Anner")
LOGIN_PASSWORD = os.environ.get("LOGIN_PASSWORD", "1")
POLL_SECONDS = int(os.environ.get("POLL_SECONDS", "240"))
PG_SSH_HOST = os.environ.get("PG_SSH_HOST", "documind")
PG_CONTAINER = os.environ.get("PG_CONTAINER", "documind-postgres")
PG_USER = os.environ.get("PG_USER", "documind")
PG_DATABASE = os.environ.get("PG_DATABASE", "documind_dev")


def fail(message, details=None):
    print(f"FAIL: {message}")
    if details is not None:
        print(json.dumps(details, ensure_ascii=False, indent=2))
    raise SystemExit(1)


def ok(message):
    print(f"PASS: {message}")


def require_optional_modules():
    missing = []
    for module in ("PIL", "reportlab", "pptx"):
        try:
            __import__(module)
        except Exception:
            missing.append(module)
    if missing:
        fail(
            "missing Python modules required for preview/OCR smoke",
            {"missing": missing, "hint": "install Pillow reportlab python-pptx"},
        )


def http_json(method, path, payload=None, token=None, timeout=60):
    data = None
    headers = {"Accept": "application/json"}
    if payload is not None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        headers["Content-Type"] = "application/json"
    if token:
        headers["Authorization"] = f"Bearer {token}"
    req = urllib.request.Request(
        f"{BASE_URL}{path}", data=data, headers=headers, method=method
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            body = resp.read().decode("utf-8")
            return json.loads(body) if body else None
    except urllib.error.HTTPError as err:
        body = err.read().decode("utf-8", "replace")
        try:
            parsed = json.loads(body)
        except json.JSONDecodeError:
            parsed = body
        fail(f"{method} {path} returned HTTP {err.code}", parsed)


def http_bytes(path, token=None, timeout=120):
    headers = {"Accept": "*/*"}
    if token:
        headers["Authorization"] = f"Bearer {token}"
    req = urllib.request.Request(f"{BASE_URL}{path}", headers=headers, method="GET")
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read(), resp.headers.get("Content-Type", "")


def http_json_url(url, timeout=120):
    req = urllib.request.Request(f"{BASE_URL}{url}", headers={"Accept": "application/json"}, method="GET")
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return json.loads(resp.read().decode("utf-8"))


def http_bytes_url(url, timeout=120):
    req = urllib.request.Request(f"{BASE_URL}{url}", headers={"Accept": "*/*"}, method="GET")
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read(), resp.headers.get("Content-Type", "")


def sse_post(path, payload, token):
    req = urllib.request.Request(
        f"{BASE_URL}{path}",
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Accept": "text/event-stream",
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    events = []
    current_event = None
    data_lines = []
    with urllib.request.urlopen(req, timeout=240) as resp:
        for raw_line in resp:
            line = raw_line.decode("utf-8").rstrip("\n")
            if not line:
                if current_event:
                    data_text = "\n".join(data_lines)
                    data = json.loads(data_text) if data_text else {}
                    events.append({"event": current_event, "data": data})
                    if current_event in {"answer.completed", "answer.failed"}:
                        break
                current_event = None
                data_lines = []
                continue
            if line.startswith("event:"):
                current_event = line.split(":", 1)[1].strip()
            elif line.startswith("data:"):
                data_lines.append(line.split(":", 1)[1].strip())
    return events


def answer_text(events):
    return "".join(
        event["data"].get("text", "")
        for event in events
        if event["event"] == "answer.delta"
    )


def citations_from_events(events):
    citations = []
    for event in events:
        if event["event"] == "citation.delta":
            citation = event["data"].get("citation")
            if citation:
                citations.append(citation)
    return citations


def xml_escape(text):
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def write_docx(path, paragraphs):
    body = "".join(
        f"<w:p><w:r><w:t>{xml_escape(paragraph)}</w:t></w:r></w:p>"
        for paragraph in paragraphs
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        zf.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        zf.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>{body}<w:sectPr/></w:body>
</w:document>""",
        )


def write_pptx(path, title, lines):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(914400, 1828800, 7315200, 2743200)
    frame = textbox.text_frame
    frame.text = lines[0]
    for line in lines[1:]:
        para = frame.add_paragraph()
        para.text = line
        para.level = 0
    prs.save(path)


def write_scanned_pdf(path, marker, code, city):
    from PIL import Image, ImageDraw, ImageFont
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    image = Image.new("RGB", (1800, 900), "white")
    draw = ImageDraw.Draw(image)
    font = None
    for font_path in (
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    ):
        if Path(font_path).exists():
            font = ImageFont.truetype(font_path, 72)
            break
    if font is None:
        font = ImageFont.load_default()
    lines = [
        "DocuMind OCR smoke document",
        f"Marker: {marker}",
        f"Verification code: {code}",
        f"City: {city}",
    ]
    y = 120
    for line in lines:
        draw.text((120, y), line, fill="black", font=font)
        y += 130

    pdf = canvas.Canvas(str(path), pagesize=letter)
    pdf.drawImage(ImageReader(image), 36, 250, width=540, height=270)
    pdf.showPage()
    pdf.save()


def content_type_for(path):
    guessed, _ = mimetypes.guess_type(path.name)
    return guessed or "application/octet-stream"


def multipart_upload(path, file_path, token):
    boundary = f"----documind{uuid.uuid4().hex}"
    content = file_path.read_bytes()
    body = b"".join(
        [
            f"--{boundary}\r\n".encode(),
            (
                f'Content-Disposition: form-data; name="file"; filename="{file_path.name}"\r\n'
                f"Content-Type: {content_type_for(file_path)}\r\n\r\n"
            ).encode(),
            content,
            b"\r\n",
            f"--{boundary}--\r\n".encode(),
        ]
    )
    req = urllib.request.Request(
        f"{BASE_URL}{path}",
        data=body,
        headers={
            "Accept": "application/json",
            "Authorization": f"Bearer {token}",
            "Content-Type": f"multipart/form-data; boundary={boundary}",
        },
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=120) as resp:
        return json.loads(resp.read().decode("utf-8"))


def poll_document(doc_id, token, accepted_statuses):
    deadline = time.time() + POLL_SECONDS
    last = None
    while time.time() < deadline:
        detail = http_json("GET", f"/api/admin/documents/{doc_id}", token=token)
        doc = detail.get("document") or {}
        status = doc.get("parse_status")
        if status != last:
            print(f"document {doc_id} status={status} chunks={doc.get('chunk_count')}")
            last = status
        if status in accepted_statuses:
            return detail
        if status in {"parse_failed", "embedding_failed"}:
            fail(f"document {doc_id} failed", detail)
        time.sleep(3)
    fail(f"document {doc_id} did not reach {sorted(accepted_statuses)} in time")


def run_pg_query(sql):
    if not PG_SSH_HOST:
        fail("PG_SSH_HOST is required for OCR DB validation")
    container = shlex.quote(PG_CONTAINER)
    user = shlex.quote(PG_USER)
    database = shlex.quote(PG_DATABASE)
    remote = (
        f"podman exec -i {container} psql -U {user} -d {database} "
        "-qAtX -v ON_ERROR_STOP=1"
    )
    return subprocess.run(
        ["ssh", PG_SSH_HOST, remote],
        input=f"SET search_path TO documind, public;\n{sql}",
        text=True,
        capture_output=True,
        check=True,
    ).stdout


def login():
    payload = {"email": LOGIN_EMAIL, "password": LOGIN_PASSWORD}
    data = http_json("POST", "/api/v1/auth/login", payload=payload)
    token = data.get("access_token")
    kb_ids = data.get("allowed_kb_ids") or []
    if not token or not kb_ids:
        fail("login response missing token or allowed knowledge base", data)
    ok("login succeeded")
    return token, kb_ids[0]


def assert_pdf_bytes(label, data, content_type):
    if not data.startswith(b"%PDF"):
        fail(f"{label} did not return PDF bytes", {"content_type": content_type})
    ok(f"{label} returned PDF bytes")


def verify_office_preview(doc_id, token, expected_format):
    preview = http_json("GET", f"/api/files/{doc_id}/preview", token=token)
    if preview.get("preview_type") != "office_pdf":
        fail("office preview type mismatch", preview)
    manifest = http_json("GET", f"/api/files/{doc_id}/preview/manifest", token=token, timeout=180)
    if manifest.get("format") != expected_format:
        fail("office manifest format mismatch", manifest)
    if manifest.get("preview_type") != "office_pdf":
        fail("office manifest preview_type mismatch", manifest)
    if manifest.get("conversion_status") != "converted":
        fail("office preview was not converted", manifest)
    if not manifest.get("page_count") or manifest.get("page_count") < 1:
        fail("office preview page_count missing", manifest)
    data, content_type = http_bytes(f"/api/files/{doc_id}/preview/content", token=token)
    assert_pdf_bytes(f"{expected_format} preview content", data, content_type)
    page, page_type = http_bytes(f"/api/files/{doc_id}/preview/pages/1/pdf", token=token)
    assert_pdf_bytes(f"{expected_format} preview page", page, page_type)
    signed = http_json("GET", f"/api/files/{doc_id}/preview-url", token=token)
    for field in ("expires_at", "preview_url", "manifest_url", "page_pdf_url_template"):
        if not signed.get(field):
            fail("preview-url response missing field", signed)
    signed_manifest = http_json_url(signed["manifest_url"], timeout=180)
    if signed_manifest.get("preview_type") != "office_pdf":
        fail("signed preview manifest mismatch", signed_manifest)
    signed_content, signed_content_type = http_bytes_url(signed["preview_url"])
    assert_pdf_bytes(f"{expected_format} signed preview content", signed_content, signed_content_type)
    signed_page_url = signed["page_pdf_url_template"].replace("{page}", "1")
    signed_page, signed_page_type = http_bytes_url(signed_page_url)
    assert_pdf_bytes(f"{expected_format} signed preview page", signed_page, signed_page_type)
    ok(f"{expected_format} office preview manifest is converted")


def upload_and_index(path, token, kb_id):
    uploaded = multipart_upload(f"/api/knowledge-bases/{kb_id}/documents", path, token)
    doc_id = uploaded.get("document_id")
    if not doc_id:
        fail("upload response missing document_id", uploaded)
    poll_document(doc_id, token, {"indexed"})
    return doc_id


def verify_ocr(token, kb_id, tmpdir, marker):
    code = "73941"
    city = "HANGZHOU"
    path = Path(tmpdir) / f"ocr-smoke-{marker}.pdf"
    write_scanned_pdf(path, marker, code, city)
    uploaded = multipart_upload(f"/api/knowledge-bases/{kb_id}/documents", path, token)
    doc_id = uploaded.get("document_id")
    if not doc_id:
        fail("OCR upload response missing document_id", uploaded)
    poll_document(doc_id, token, {"parse_low_confidence"})
    queued = http_json("POST", f"/api/admin/documents/{doc_id}/send-to-ocr", token=token)
    if queued.get("parse_status") != "ocr_pending":
        fail("send-to-ocr did not queue OCR", queued)
    poll_document(doc_id, token, {"indexed"})

    rows = run_pg_query(
        f"""
SELECT
  d.parse_status,
  d.metadata->>'ocr_status',
  d.chunk_count,
  (SELECT count(*) FROM chunks c WHERE c.doc_id = d.id AND c.content ILIKE '%{marker}%'),
  (SELECT count(*) FROM document_source_anchors a WHERE a.doc_id = d.id AND a.bbox IS NOT NULL)
FROM documents d
WHERE d.id = '{doc_id}';
"""
    ).strip()
    if not rows:
        fail("OCR document missing from PostgreSQL", {"doc_id": doc_id})
    status, ocr_status, chunk_count, marker_chunks, bbox_anchors = rows.split("|")
    if status != "indexed" or ocr_status != "completed":
        fail("OCR metadata did not complete", rows)
    if int(chunk_count) < 1 or int(marker_chunks) < 1 or int(bbox_anchors) < 1:
        fail("OCR chunks or bbox anchors missing", rows)
    ok("OCR DB outputs include chunks, marker text, and bbox anchors")

    conv = http_json(
        "POST",
        "/api/conversations",
        {"title": f"OCR smoke {marker}", "kb_ids": [kb_id]},
        token=token,
    )
    conv_id = conv["conversation_id"]
    events = sse_post(
        f"/api/conversations/{conv_id}/messages",
        {
            "content": f"在 OCR smoke 文档 {marker} 中，验证码和城市分别是什么？",
            "kb_ids": [],
            "client_request_id": f"ocr-smoke-{uuid.uuid4()}",
            "stream": True,
        },
        token,
    )
    answer = answer_text(events)
    citations = citations_from_events(events)
    normalized = answer.lower()
    if code not in answer or city.lower() not in normalized:
        fail("OCR answer did not contain expected code and city", {"answer": answer})
    if not any(c.get("doc_id") == doc_id and c.get("anchor") for c in citations):
        fail("OCR answer missing citation anchor for OCR document", citations)
    ok("OCR QA returned expected answer and citation anchor")


def main():
    require_optional_modules()
    token, kb_id = login()
    marker = f"PVSMOKE-{uuid.uuid4().hex[:8].upper()}"
    with tempfile.TemporaryDirectory(prefix="documind-preview-ocr-") as tmpdir:
        docx = Path(tmpdir) / f"office-docx-{marker}.docx"
        write_docx(
            docx,
            [
                f"DocuMind preview smoke {marker}",
                "This DOCX verifies LibreOffice PDF preview conversion.",
            ],
        )
        pptx = Path(tmpdir) / f"office-pptx-{marker}.pptx"
        write_pptx(
            pptx,
            f"DocuMind preview smoke {marker}",
            ["This PPTX verifies Office PDF preview conversion.", marker],
        )

        docx_id = upload_and_index(docx, token, kb_id)
        verify_office_preview(docx_id, token, "docx")
        pptx_id = upload_and_index(pptx, token, kb_id)
        verify_office_preview(pptx_id, token, "pptx")
        verify_ocr(token, kb_id, tmpdir, marker)

    print("preview/OCR smoke passed")


if __name__ == "__main__":
    main()
